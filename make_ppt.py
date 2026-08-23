import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
blank_layout = prs.slide_layouts[6]

C_DARK_BG = RGBColor(0x08, 0x12, 0x1E)
C_ACCENT1 = RGBColor(0x00, 0xD4, 0xFF)
C_ACCENT2 = RGBColor(0x00, 0xB4, 0x8A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
C_BOX = RGBColor(0x16, 0x2A, 0x40)

def bg(slide): slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C_DARK_BG
def txbox(slide, text, left, top, width, height, size=16, bold=False, color=C_WHITE):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text, r.font.size, r.font.bold, r.font.color.rgb = text, Pt(size), bold, color
    return tb

slides_data = [
    ("MediChain Enterprise v3.0", "Blockchain-Based EHR with AI-Assisted Clinical Decision Intelligence"),
    ("Problem Statement", "Data Fragmentation, Privacy Breaches, No Predictive Care"),
    ("Architecture", "4-Tier: React SPA, Express API, Smart Contract (Ethereum), Python Flask AI"),
    ("Blockchain Layer", "MediChain.sol, IPFS storage (Pinata), AES-256-GCM encryption, Tamper Proofing"),
    ("Backend Layer", "Node.js, Express, MongoDB Atlas, JWT Auth, Role-Based Access Control"),
    ("AI / CDSS", "6 Disease Models, Drug Interaction Engine, Dosage Safety ML, Predictive Analytics"),
    ("Frontend SPA", "React 18, TailwindCSS, 28 Dashboards, QR Health ID, Health Timeline"),
    ("Security & Privacy", "OWASP-Compliant, Patient-controlled access, Time-limited doctor grants"),
    ("Conclusion", "100% Tests Passing, Production Ready, Multi-Layer Healthcare Platform")
]

for title, body in slides_data:
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    txbox(s, title, 1, 1, 11, 1, size=40, bold=True, color=C_ACCENT1)
    txbox(s, body, 1, 3, 11, 3, size=24, color=C_LIGHT_GRAY)

prs.save("MediChain_v3_Presentation.pptx")
print("Presentation generated: MediChain_v3_Presentation.pptx")
