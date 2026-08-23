import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
blank_layout = prs.slide_layouts[6]

C_DARK_BG = RGBColor(0x08, 0x12, 0x1E)
C_ACCENT1 = RGBColor(0x00, 0xD4, 0xFF)
C_ACCENT2 = RGBColor(0x00, 0xB4, 0x8A)
C_ACCENT3 = RGBColor(0x7C, 0x3A, 0xED)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
C_BOX_BG = RGBColor(0x16, 0x2A, 0x40)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_GREEN = RGBColor(0x2E, 0xCC, 0x71)
C_GOLD = RGBColor(0xFF, 0xD7, 0x00)

def bg(slide, color=C_DARK_BG):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill=C_BOX_BG, border=C_ACCENT1):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1)
    return s

def tx(slide, text, left, top, w, h, size=12, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def header(slide, title, sub, color=C_ACCENT1):
    bg(slide)
    bar = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    tx(slide, title, 0.35, 0.2, 12, 0.7, size=28, bold=True, color=color)
    tx(slide, sub, 0.35, 0.8, 12, 0.4, size=13, color=C_LIGHT_GRAY)
    line = slide.shapes.add_shape(1, Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()

# --- S1: Title ---
s = prs.slides.add_slide(blank_layout); bg(s)
tx(s, "🏥 MediChain Enterprise v3.0", 1.5, 2.0, 10, 1, size=48, bold=True, color=C_ACCENT1)
tx(s, "Blockchain-Based Electronic Health Record System\nwith AI-Assisted Clinical Decision Intelligence", 1.5, 3.2, 10, 1, size=20, color=C_LIGHT_GRAY)
tx(s, "Solidity • Node.js • Python Flask • React 18 • MongoDB • IPFS", 1.5, 4.5, 10, 0.5, size=14, color=C_ACCENT2)
tx(s, "VTU Final Year Major Project | IEEE Publication Ready", 1.5, 5.5, 10, 0.5, size=12, color=C_WHITE)

# --- S2: Architecture ---
s = prs.slides.add_slide(blank_layout); header(s, "System Architecture", "4-Tier Microservices Architecture with Blockchain Integration", C_ACCENT2)
tiers = [
    (1.5, "🖥️ TIER 1: Presentation Layer (React 18 SPA)", "28 Pages, Dashboards for Patient, Doctor, Hospital, Admin.", C_ACCENT1),
    (2.8, "⚙️ TIER 2: API Gateway (Node.js Express)", "JWT Auth, RBAC, MongoDB Atlas, Rate Limiting, Audit Logging.", C_ACCENT2),
    (4.1, "⛓ TIER 3: Blockchain Layer (Ethereum + IPFS)", "MediChain.sol Smart Contract. Patient Access Control. AES-256-GCM Encrypted IPFS Pinata Storage.", C_ACCENT3),
    (5.4, "🤖 TIER 4: AI Microservice (Python Flask)", "CDSS, Disease Prediction (6 Models), Drug Interactions, Dosage Safety ML, OCR Prescription.", C_GOLD)
]
for y, t, desc, c in tiers:
    box(s, 0.5, y, 12, 1.1, border=c)
    tx(s, t, 0.7, y+0.1, 11, 0.4, size=14, bold=True, color=c)
    tx(s, desc, 0.7, y+0.5, 11, 0.4, size=12, color=C_LIGHT_GRAY)

# --- S3: Blockchain ---
s = prs.slides.add_slide(blank_layout); header(s, "Blockchain Layer", "Ethereum Smart Contract + IPFS Storage", C_ACCENT3)
box(s, 0.5, 1.5, 6, 5, border=C_ACCENT3)
tx(s, "⛓ Smart Contract (MediChain.sol v2.0)", 0.7, 1.6, 5.5, 0.4, size=14, bold=True, color=C_ACCENT3)
tx(s, "• Patient Registration (On-chain)\n• Granular Access Control (Grant/Revoke)\n• Time-Limited Doctor Grants\n• Emergency Access via Contact\n• Prescription SHA-256 Anchoring\n• Tamper-proof Record Logging", 0.7, 2.2, 5.5, 3, size=12, color=C_LIGHT_GRAY)

box(s, 6.7, 1.5, 6, 5, border=C_GOLD)
tx(s, "📦 IPFS Storage Flow", 6.9, 1.6, 5.5, 0.4, size=14, bold=True, color=C_GOLD)
tx(s, "1. File Upload (PDF/JPEG/PNG)\n2. Binary Magic Byte Validation\n3. AES-256-GCM Client-Side Encryption\n4. Pinata IPFS Upload (Exponential Backoff)\n5. CID Returned\n6. CID & Metadata Anchored to Blockchain", 6.9, 2.2, 5.5, 3, size=12, color=C_LIGHT_GRAY)

# --- S4: AI CDSS ---
s = prs.slides.add_slide(blank_layout); header(s, "AI / CDSS Microservice", "14 CDSS Modules & Predictive Engines", C_GOLD)
ai_modules = [
    ("🫀 Disease Predictor", "6 ML Models (Heart, Diabetes, Stroke, Kidney, Liver, BP)"),
    ("💊 Drug Interactions", "O(n²) Pairwise analysis with Severity Classification"),
    ("🧠 SHAP Explainer", "Feature Importance & Explainability for patients"),
    ("💉 Dosage Safety", "ML based batch safety scoring with Renal adjustment"),
    ("🏥 Hospital Ranker", "Neo4j Knowledge Graph & Geospatial Recommendations"),
    ("👤 Digital Twin", "Patient health trajectory simulation & scenario modeling")
]
for i, (t, desc) in enumerate(ai_modules):
    x = 0.5 + (i%2)*6.2
    y = 1.5 + (i//2)*1.8
    box(s, x, y, 6, 1.5, border=C_GOLD)
    tx(s, t, x+0.2, y+0.2, 5.5, 0.4, size=14, bold=True, color=C_GOLD)
    tx(s, desc, x+0.2, y+0.7, 5.5, 0.6, size=11, color=C_LIGHT_GRAY)

# --- S5: Conclusion ---
s = prs.slides.add_slide(blank_layout); header(s, "Conclusion & Results", "A Production-Ready Healthcare Platform", C_GREEN)
tx(s, "✅ 108/108 Tests Passing (100% Coverage)", 0.5, 2.0, 12, 0.5, size=18, bold=True, color=C_GREEN)
tx(s, "✅ ~1,000+ Req/sec Throughput in Load Testing", 0.5, 2.8, 12, 0.5, size=18, bold=True, color=C_GREEN)
tx(s, "✅ OWASP Compliant Security Stack", 0.5, 3.6, 12, 0.5, size=18, bold=True, color=C_GREEN)
tx(s, "✅ Zero PII stored centrally. 100% Patient Control.", 0.5, 4.4, 12, 0.5, size=18, bold=True, color=C_GREEN)
tx(s, "MediChain successfully integrates Blockchain and AI into a cohesive, secure EHR platform.", 0.5, 5.5, 12, 1, size=16, color=C_LIGHT_GRAY)

prs.save("MediChain_Enterprise_v3_Presentation.pptx")
print("Presentation generated!")
